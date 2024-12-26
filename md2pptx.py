# md2pptx.py
import os
import re
import datetime
from markdown import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from dotenv import dotenv_values

def preprocess_marp(md_text):
    md_text = re.sub(r'<!--.*?-->', '', md_text, flags=re.DOTALL)
    return md_text.strip()

def parse_front_matter(md_text):
    fm_pattern = r'^---\s*(.*?)\s*---\s*(.*)$'
    match = re.search(fm_pattern, md_text, flags=re.DOTALL|re.MULTILINE)
    front_matter = {}
    body = md_text
    if match:
        fm_text = match.group(1)
        body = match.group(2)
        for line in fm_text.split('\n'):
            line=line.strip()
            if not line:
                continue
            if ':' in line:
                key, val = line.split(':',1)
                key = key.strip()
                val = val.strip()
                front_matter[key] = val
    return front_matter, body

def split_slides(md_text):
    slides = [s.strip() for s in md_text.split('---') if s.strip()]
    return slides

def parse_slide(slide_md):
    html = markdown(slide_md, extensions=['extra', 'smarty'])
    soup = BeautifulSoup(html, 'html.parser')
    notes_div = soup.find('div', class_='notes')
    notes = None
    if notes_div:
        notes = notes_div.get_text(strip=True, separator='\n')
        notes_div.decompose()
    heading = soup.find(['h1', 'h2', 'h3'])
    title = heading.get_text(strip=True) if heading else None
    if heading:
        heading.decompose()

    bullets = [li.get_text(strip=True) for li in soup.find_all('li')]
    paragraphs = [p.get_text(strip=True) for p in soup.find_all('p')]

    if not title:
        if paragraphs:
            title = paragraphs[0]
            paragraphs = paragraphs[1:]
        else:
            title = "No Title"
    if bullets and not paragraphs:
        slide_type = 'bullet'
        text_lines = []
    elif paragraphs and not bullets:
        slide_type = 'text'
        text_lines = paragraphs
    elif bullets and paragraphs:
        slide_type = 'bullet'
        bullets.extend(paragraphs)
        text_lines = []
    else:
        slide_type = 'text'
        text_lines = []

    slide_data = {
        'type': slide_type,
        'title': title,
        'notes': notes,
    }
    if slide_type == 'bullet':
        slide_data['bullets'] = bullets
    elif slide_type == 'text':
        slide_data['text_lines'] = text_lines

    return slide_data

def build_slides_data(md_text):
    slides = split_slides(md_text)
    slides_data = [parse_slide(s) for s in slides]
    return slides_data

def load_size_map():
    env = dotenv_values(".env")
    size_dict = {}
    for k,v in env.items():
        if k.startswith("SIZE_MAP_"):
            ratio_key = k.replace("SIZE_MAP_", "").replace("_", ":")
            w,h = v.split(',')
            w = float(w.strip())
            h = float(h.strip())
            size_dict[ratio_key] = (Inches(w), Inches(h))
    return size_dict

def add_title_slide(prs, title_text, subtitle_text=None, notes=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    if subtitle_text is not None and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_bullet_slide(prs, title_text, bullet_points, title_layout=1, font_size=18, notes=None):
    slide_layout = prs.slide_layouts[title_layout]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bp in bullet_points:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0
        p.font.size = Pt(font_size)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_text_slide(prs, title_text, text_lines, title_layout=1, font_size=18, bold=False, notes=None):
    slide_layout = prs.slide_layouts[title_layout]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for line in text_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def create_ppt_from_slides_data(front_matter, slides_data, output_filename="output.pptx", template_file=None):
    if template_file:
        prs = Presentation(template_file)
    else:
        prs = Presentation()
    size_map = load_size_map()
    slide_ratio = front_matter.get('size', None)
    DEFAULT_SIZE = (Inches(10), Inches(7.5))
    if slide_ratio in size_map:
        prs.slide_width, prs.slide_height = size_map[slide_ratio]
    else:
        prs.slide_width, prs.slide_height = DEFAULT_SIZE

    title_text = front_matter.get('title', 'No Title')
    now = datetime.datetime.now()
    today = now.strftime('%Y年%m月%d日')
    add_title_slide(prs, title_text, subtitle_text=today)

    for slide_data in slides_data:
        if slide_data['type'] == 'bullet':
            add_bullet_slide(prs, slide_data['title'], slide_data['bullets'], notes=slide_data['notes'])
        elif slide_data['type'] == 'text':
            add_text_slide(prs, slide_data['title'], slide_data.get('text_lines', []), notes=slide_data['notes'])
        else:
            add_text_slide(prs, slide_data['title'], [], notes=slide_data['notes'])

    prs.save(output_filename)
    print(f"PowerPointファイル出力完了: {output_filename}")

def convert_md_to_pptx(md_file, pptx_output_file="output.pptx"):
    with open(md_file, 'r', encoding='utf-8') as f:
        md_text = f.read()
    md_text = preprocess_marp(md_text)
    front_matter, body = parse_front_matter(md_text)
    slides_data = build_slides_data(body)
    create_ppt_from_slides_data(front_matter, slides_data, pptx_output_file)